import openai
import streamlit as st
import io
from docx import Document
from pptx import Presentation
import pandas as pd
import PyPDF2
from pdf2docx import Converter
import tempfile
import os

# OpenAI 클라이언트 생성
openai.api_key = st.secrets["OPENAI_API_KEY"]
# 번역 스타일별 예시 확장
parallel_example = {
    "한국어": {
        "Academic": [
            "인공 지능 (AI) 은 컴퓨터 과학 내에서 새롭게 떠오르고 빠르게 진화하는 영역으로, 인간이 전통적으로 수행하던 작업을 수행할 수 있는 지능형 기계 개발에 전념하고 있습니다.의료, 금융, 교통 등 다양한 분야에서 AI 기술의 확산이 점점 더 뚜렷해지고 있습니다.효율성, 정확성 및 의사 결정 프로세스를 개선할 수 있는 AI의 잠재력은 수많은 산업을 근본적으로 변화시킬 가능성을 내포하고 있습니다 AI 시스템은 머신러닝 알고리즘의 적용을 통해 지식을 습득하고 기능을 개선합니다.이러한 알고리즘은 로봇 시스템이 광범위한 데이터 세트를 분석할 수 있도록 하여 사람이 감지하지 못할 수 있는 패턴과 통찰력을 발견할 수 있도록 합니다.이러한 발전은 자연어 처리, 컴퓨터 비전 및 음성 인식 분야를 크게 발전시켰습니다.그럼에도 불구하고 AI의 도입은 잠재적 직업 이직 및 의사 결정의 편향에 대한 우려를 불러일으키고 있습니다.따라서 AI 연구 및 개발에는 윤리적 고려 사항이 반드시 포함되어야 합니다 AI는 사회 복지를 증진하고 삶의 질을 개선할 수 있는 역량을 가지고 있지만, 그 적용에는 윤리 원칙이 적용되어야 합니다."
        ],
        "Fluent": [
            "인공 지능 (AI) 은 일반적으로 인간이 수행하는 작업을 실행할 수 있는 지능형 시스템을 고안하는 컴퓨터 과학 내에서 빠르게 발전하는 분야입니다.의료, 금융 및 교통과 같은 분야에서 AI 기술의 적용이 점점 더 보편화되고 있습니다 AI의 효과, 정확성 및 의사 결정 능력은 수많은 산업에 혁명을 일으킬 잠재력을 가지고 있습니다 AI는 기계 학습 알고리즘을 통해 기능을 향상시키고 진화합니다.이러한 알고리즘을 통해 머신은 방대한 양의 데이터를 분석하고 인간의 인식을 벗어나는 패턴과 통찰력을 발견할 수 있습니다.이러한 발전으로 자연어 처리, 컴퓨터 비전 및 음성 인식이 크게 발전했습니다.그럼에도 불구하고 AI는 직업 이동 및 의사 결정 과정의 편향과 같은 위험도 내포하고 있습니다.따라서 AI 연구 및 개발에는 윤리적 고려 사항을 반영하는 것이 필수적입니다 AI는 우리 사회를 개선하고 삶을 개선할 수 있다는 가능성을 지니고 있지만, 윤리적 관행에 기반을 두고 이를 구현하는 것이 매우 중요합니다."
        ],
        "Formal": [
            "인공 지능 (AI) 은 전통적으로 인간이 수행하던 작업을 수행할 수 있는 지능형 기계 개발에 전념하는 컴퓨터 과학 내에서 빠르게 발전하는 분야입니다 AI 기술의 확산은 의료, 은행 및 교통과 같은 분야에서 점점 더 분명해지고 있습니다.효율성, 정확성, 의사 결정 능력 등 AI의 역량은 수많은 산업에 혁명을 일으킬 잠재력을 가지고 있습니다 AI는 머신러닝 알고리즘의 적용을 통해 기능과 성능을 향상시킵니다.이러한 알고리즘은 기계가 광범위한 데이터 세트를 평가하고 인간의 이해를 넘어서는 패턴과 통찰력을 식별할 수 있도록 합니다.이러한 발전은 자연어 처리, 컴퓨터 비전 및 음성 인식 분야를 크게 발전시켰습니다.그럼에도 불구하고 AI의 출현은 직업 이동과 의사 결정 과정의 편향으로 이어질 수 있습니다.따라서 AI 연구 및 개발에는 윤리적 고려 사항이 반드시 포함되어야 합니다 AI는 사회를 개선하고 삶의 질을 개선할 수 있는 역량을 갖추고 있지만, 이를 배포하려면 윤리적 기준을 반드시 준수해야 합니다."
        ],
        "Creative": [
            "인공 지능 (AI) 은 일반적으로 인간에게만 맡겨지는 작업을 수행할 수 있는 영리한 기계를 만드는 컴퓨터 과학 분야에서 매우 흥미롭게 발전하고 있습니다 AI 기술의 채택은 의료, 은행, 교통과 같은 분야에 빠르게 스며들고 있습니다 AI의 놀라운 효율성, 정확성 및 의사 결정 능력은 수많은 산업에 혁명을 일으킬 잠재력을 가지고 있습니다.머신 러닝 알고리즘의 마법을 통해 AI는 지속적으로 학습하고 기능을 개선합니다.이러한 알고리즘은 로봇이 방대한 데이터 세트를 샅샅이 뒤져 사람이 파악할 수 없는 패턴과 통찰력을 찾아낼 수 있도록 합니다.이러한 발전은 자연어 처리, 컴퓨터 비전 및 음성 인식 분야의 발전을 새로운 차원으로 끌어올렸습니다.그럼에도 불구하고 AI는 직업 이동 및 의사 결정 과정의 편향과 같은 문제를 야기합니다.따라서 AI 연구 및 개발 노력의 최전선에 윤리적 고려가 있어야 합니다 AI는 우리 사회를 개선하고 삶을 풍요롭게 할 수 있는 가능성을 지니고 있지만, AI의 구현은 윤리 원칙에 따라 이루어져야 합니다."
        ],
        "Informal": [
            "인공 지능 (AI) 은 컴퓨터 과학의 매우 트렌디한 분야와 같습니다. 인간이 할 수 있는 일을 할 수 있는 기계를 만드는 것이 전부입니다 AI는 어디에나 나타나고 있습니다. 특히 의료, 은행, 심지어 도시 곳곳을 돌아다니는 곳에서요 AI의 작동 방식은 속도와 현명한 의사 결정으로 많은 산업을 완전히 뒤흔들 수 있습니다.수많은 데이터를 선별하고 우리가 놓칠 수 있는 패턴을 찾아내는 데 도움이 되는 기계 학습 알고리즘 덕분에 더 나은 성과를 거둘 수 있습니다.이 기술은 언어를 이해하고, 사물을 보고, 음성을 인식하는 것과 같은 분야에서 큰 발전을 이루었습니다.하지만 햇빛과 무지개만 있는 것은 아닙니다. AI는 일자리를 잃고 개략적인 의사 결정으로 이어질 수도 있습니다.따라서 AI를 연구할 때는 관련된 윤리에 대해 생각해 보는 것이 매우 중요합니다 AI는 우리의 삶과 사회를 더 좋게 만들 수 있는 잠재력을 가지고 있습니다. 하지만 우리는 AI를 올바른 방법으로 사용해야 합니다"
        ],
    },
    "영어": {
        "Academic": [
            "Artificial Intelligence (AI) represents a nascent and swiftly advancing sector within the discipline of computer science, focused on the creation of intelligent machines that are proficient at performing tasks conventionally linked to human cognitive abilities. The burgeoning presence of AI technologies is increasingly observable across various sectors, including healthcare, finance, and transportation. The promise of AI to augment efficiency, precision, and decision-making processes heralds the potential to fundamentally transform a wide array of industries. AI systems acquire knowledge and improve their functionalities through intricate machine learning algorithms. These algorithms enable automated systems to scrutinize vast datasets and identify patterns and insights that remain beyond human recognition. This advancement has substantially progressed domains such as natural language processing, computer vision, and speech recognition. However, the deployment of AI engenders apprehensions concerning possible job displacement and biases within decision-making frameworks. Consequently, it is essential that the research and development of AI integrate ethical considerations. While AI possesses the potential to elevate societal welfare and enhance quality of life, it is of utmost importance that its applications conform to ethical standards."
        ],
        "Fluent": [
            "Artificial Intelligence (AI) represents a burgeoning and swiftly advancing sector within computer science, focused on creating intelligent machines that can perform tasks generally linked to human thought processes. The rise of AI technologies is becoming increasingly apparent in various fields such as healthcare, finance, and transportation. The ability of AI to improve efficiency, precision, and decision-making capabilities carries the potential to fundamentally transform numerous industries. AI systems gain knowledge and improve their functionalities through advanced machine learning algorithms. These algorithms enable automated systems to scrutinize large datasets and identify patterns and insights that may escape human awareness. This advancement has significantly propelled fields like natural language processing, computer vision, and speech recognition. However, the deployment of AI raises issues regarding possible job displacement and biases within decision-making processes. Consequently, it is essential for the research and development of AI to integrate ethical considerations. While AI has the potential to enhance societal welfare and elevate quality of life, it is vital that its applications conform to ethical norms."
        ],
        "Formal": [
            "Artificial Intelligence (AI) represents a nascent and swiftly advancing sector within the discipline of computer science, focused on the creation of intelligent machines that are proficient at performing tasks conventionally linked to human cognitive abilities. The burgeoning presence of AI technologies is increasingly observable across various sectors, including healthcare, finance, and transportation. The promise of AI to augment efficiency, precision, and decision-making processes heralds the potential to fundamentally transform a wide array of industries. AI systems acquire knowledge and improve their functionalities through intricate machine learning algorithms. These algorithms enable automated systems to scrutinize vast datasets and identify patterns and insights that remain beyond human recognition. This advancement has substantially progressed domains such as natural language processing, computer vision, and speech recognition. However, the deployment of AI engenders apprehensions concerning possible job displacement and biases within decision-making frameworks. Consequently, it is essential that the research and development of AI integrate ethical considerations. While AI possesses the potential to elevate societal welfare and enhance quality of life, it is of utmost importance that its applications conform to ethical standards."
        ],
        "Creative": [
            "Artificial Intelligence (AI) represents a groundbreaking and swiftly transforming frontier within the realm of computer science, committed to crafting intelligent machines that can perform tasks traditionally linked to human thought processes. The surge of AI innovations is unmistakably manifesting across diverse fields like healthcare, finance, and transportation. The ability of AI to amplify efficiency, precision, and decision-making capabilities heralds the potential to revolutionize numerous industries. AI systems gain wisdom and bolster their skills through advanced machine learning frameworks. These frameworks enable automated systems to sift through vast troves of data and uncover patterns and insights that often escape human awareness. This evolution has markedly propelled areas like natural language processing, computer vision, and speech recognition into new heights. However, the deployment of AI brings forth apprehensions about possible job losses and biases in the decision-making landscape. Thus, it becomes crucial that the exploration and advancement of AI embed ethical considerations at their core. While AI has the potential to elevate societal welfare and enrich lives, it is essential that its utilization remains aligned with ethical principles."
        ],
        "Informal": [
            "Artificial Intelligence (AI) is like this super cool and fast-growing thing in computer science that's all about making smart machines that can do stuff we usually think only humans can do. You can totally see AI popping up everywhere, like in healthcare, finance, and transportation. The way AI can boost how well things work, make stuff more accurate, and help with decision-making is pretty much set to change a bunch of industries for the better. AI systems learn and get better at their jobs using some fancy machine learning tricks. These tricks let automated systems dive into huge piles of data and pick out patterns and insights that we might totally miss. Because of this, we've made huge strides in areas like natural language processing, computer vision, and speech recognition. But, on the flip side, using AI also brings up worries about people losing jobs and biases in how decisions get made. So, it’s super important that we think about ethics while working on AI. Even though AI can really help society and make life better, it’s key that we keep everything above board and stick to ethical guidelines."
        ],
    },
    "중국어": {
        "Academic": [
            "人工智能（AI）是计算机科学的一个快速发展的分支，致力于开发能够执行传统上与人类能力相关的任务的智能系统。人工智能技术在医疗保健、金融和交通等不同领域的扩散越来越明显。人工智能在提高效率、准确性和决策过程方面的潜力为众多行业带来了革命性的希望。人工智能系统通过应用机器学习算法来获取知识并提高其性能。这些算法使自动化系统能够分析大量数据集并辨别出无法被人类发现的模式和见解。这一进展极大地推动了自然语言处理、计算机视觉和语音识别技术的进步。尽管如此，人工智能技术的部署带来了与工作流失和决策过程中的偏见有关的风险。因此，必须将道德考虑纳入人工智能的研发计划。尽管人工智能具有丰富我们的社会框架和改善个人生活的能力，但其实施必须遵守道德标准。"
        ],
        "Fluent": [
            "人工智能（AI）是一个快速发展的计算机科学领域，它开发出能够执行通常由人类执行的任务的智能机器。在医疗保健、银行和交通等领域，人工智能技术的使用正在显著扩展。人工智能在提高效率、准确性和决策方面的潜力可能会彻底改变众多行业。人工智能通过机器学习算法增强其能力。这些算法使机器能够分析大量数据，发现通常是人类无法察觉的模式和见解。这一进展使自然语言处理、计算机视觉和语音识别取得了重大进步。尽管如此，人工智能的兴起引发了人们对工作流失和决策过程偏见的担忧。因此，人工智能的研发必须纳入伦理方面的考虑。人工智能有望改善我们的社会，改善我们的生活，但必须以合乎道德的方式来实施。"
        ],
        "Formal": [
            "人工智能（AI）是计算机科学领域中一个快速发展的领域，致力于开发能够执行通常与人类相关的任务的智能系统。人工智能技术的应用在医疗保健、金融和交通等领域越来越普遍。人工智能在提高效率、准确性和决策过程方面的潜力可能会极大地改变各行各业。人工智能系统通过复杂的机器学习算法获取知识并完善其能力。这些算法使机器人系统能够分析大量数据集，识别出人类无法感知的模式和见解。这一进步为自然语言处理、计算机视觉和语音识别等领域的重大进展做出了贡献。尽管如此，人工智能的部署引发了人们对工作流失和决策过程偏见的担忧。因此，人工智能研发必须优先考虑伦理因素。尽管人工智能具有增强社会福利和改善生活质量的能力，但其应用必须受道德标准的约束。"
        ],
        "Creative": [
            "人工智能 (AI) 是一个令人振奋且快速发展的计算机科学领域，它造就了能够模仿人类行为的智能机器。人工智能技术的覆盖范围正在显著扩展到医疗保健、金融和交通等领域。人工智能卓越的效率、准确性和决策能力有可能彻底改变众多行业。通过机器学习算法的力量，人工智能不断学习和增强其能力。这些复杂的算法使机器能够筛选海量数据，发现人类无法感知的模式和见解。这一进展推动了自然语言处理、计算机视觉和语音识别领域的进步。尽管如此，人工智能的兴起带来了工作流失和决策过程偏见的风险。因此，人工智能的研发必须纳入伦理考量。尽管人工智能有望丰富我们的社会和改善我们的生活，但其应用必须建立在道德实践的基础上。"
        ],
        "Informal": [
            "人工智能（AI）是计算机科学中一个非常酷且发展迅速的领域，其全部目的是制造可以做人类做的事情的智能机器。现在，人工智能技术无处不在，例如医疗保健、银行甚至交通领域。人工智能如此高效、精确和擅长做出决策的方式可能会彻底改变许多行业的游戏规则。通过这些机器学习算法进行学习，人工智能会变得更好。这些算法帮助机器筛选大量数据，发现我们人类可能会错过的模式和见解。这在自然语言处理、计算机视觉和语音识别方面取得了长足的进步。但是，也有缺点；人工智能可能会导致失业和决策中的一些偏见。因此，我们在进行人工智能研发时确实需要考虑伦理问题。人工智能有可能改善我们的社会和生活，但我们必须以正确的方式使用它。"
        ],
    },
}


def translate_file(file, src_lang, trg_lang, style):
    if file.type == "application/pdf":
        # PDF를 DOCX로 변환하고 번역
        return convert_and_translate_pdf_to_docx(file, src_lang, trg_lang, style)

    elif (
        file.type
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        # DOCX 처리
        doc = Document(io.BytesIO(file.getvalue()))
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                # 원래 스타일 저장
                original_style = paragraph.style
                original_runs = []
                for run in paragraph.runs:
                    original_runs.append(
                        {
                            "text": run.text,
                            "bold": run.bold,
                            "italic": run.italic,
                            "underline": run.underline,
                            "font.name": run.font.name,
                            "font.size": run.font.size,
                            "font.color.rgb": (
                                run.font.color.rgb if run.font.color else None
                            ),
                        }
                    )

                # 텍스트 번역
                translated_text = translate_text_using_chatgpt(
                    paragraph.text, src_lang, trg_lang, style
                )

                # 번역된 텍스트로 교체하면서 스타일 유지
                paragraph.text = ""
                new_run = paragraph.add_run(translated_text)

                # 원래 스타일 적용
                if original_runs:
                    first_run = original_runs[0]
                    new_run.bold = first_run["bold"]
                    new_run.italic = first_run["italic"]
                    new_run.underline = first_run["underline"]
                    if first_run["font.name"]:
                        new_run.font.name = first_run["font.name"]
                    if first_run["font.size"]:
                        new_run.font.size = first_run["font.size"]
                    if first_run["font.color.rgb"]:
                        new_run.font.color.rgb = first_run["font.color.rgb"]
                paragraph.style = original_style

        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer

    elif (
        file.type
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        # PPTX 처리 (레이아웃 유지)
        prs = Presentation(io.BytesIO(file.getvalue()))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    # 원래 위치와 크기 저장
                    original_left = shape.left
                    original_top = shape.top
                    original_width = shape.width
                    original_height = shape.height

                    # 텍스트 번역
                    translated_text = translate_text_using_chatgpt(
                        shape.text, src_lang, trg_lang, style
                    )

                    # 번역된 텍스트로 교체
                    shape.text = translated_text

                    # 원래 위치와 크기 복원
                    shape.left = original_left
                    shape.top = original_top
                    shape.width = original_width
                    shape.height = original_height

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer

    elif (
        file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    ):
        # XLSX 처리
        df = pd.read_excel(io.BytesIO(file.getvalue()))
        translated_df = df.applymap(
            lambda x: (
                translate_text_using_chatgpt(str(x), src_lang, trg_lang, style)
                if isinstance(x, str)
                else x
            )
        )

        buffer = io.BytesIO()
        with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
            translated_df.to_excel(writer, index=False)
        buffer.seek(0)
        return buffer


def convert_and_translate_pdf_to_docx(file, src_lang, trg_lang, style):
    # 임시 파일 생성
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as pdf_temp:
        pdf_temp.write(file.getvalue())
        pdf_path = pdf_temp.name

    # PDF를 DOCX로 변환할 임시 파일 경로
    docx_path = pdf_path.replace(".pdf", ".docx")

    try:
        # PDF를 DOCX로 변환 (레이아웃 유지)
        cv = Converter(pdf_path)
        cv.convert(docx_path)
        cv.close()

        # DOCX 파일 열기
        doc = Document(docx_path)

        # 텍스트 추출 및 번역하면서 원래 스타일 유지
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():  # 빈 문단 제외
                translated_text = translate_text_using_chatgpt(
                    paragraph.text, src_lang, trg_lang, style
                )
                paragraph.text = translated_text

        # 번역된 DOCX를 메모리에 저장
        docx_buffer = io.BytesIO()
        doc.save(docx_buffer)
        docx_buffer.seek(0)

        return docx_buffer

    finally:
        # 임시 파일 삭제
        try:
            os.unlink(pdf_path)
            os.unlink(docx_path)
        except:
            pass


def get_style_instruction(style):
    instructions = {
        "Academic": {
            "instruction": "학술적이고 전문적인 어조로 번역하세요.",
            "rules": [
                "객관적이고 형식적인 어조 사용",
                "전문 용어와 학술적 표현 활용",
                "'-다' 체를 사용한 문어체",
                "인과관계와 논리적 구조를 명확히 표현",
            ],
        },
        "Fluent": {
            "instruction": "자연스럽고 부드러운 문체로 번역하세요.",
            "rules": [
                "평이하고 이해하기 쉬운 표현 사용",
                "'-입니다' 체의 격식있는 구어체",
                "전문 용어는 필요한 경우에만 사용",
                "문장 간의 자연스러운 연결",
            ],
        },
        "Formal": {
            "instruction": "격식있고 예의 바른 문체로 번역하세요.",
            "rules": [
                "공식적이고 예의 바른 표현 사용",
                "'-습니다' 체의 높임말",
                "정중하고 신중한 어조",
                "간결하고 명확한 문장 구조",
            ],
        },
        "Creative": {
            "instruction": "창의적이고 생동감 있는 문체로 번역하세요.",
            "rules": [
                "비유와 은유를 활용한 표현",
                "감정을 담은 서술적인 표현",
                "'-요' 체의 친근한 구어체",
                "독자의 흥미를 끄는 표현 사용",
            ],
        },
        "Informal": {
            "instruction": "친근하고 캐주얼한 문체로 번역하세요.",
            "rules": [
                "일상적인 대화체 사용",
                "구어적 표현과 감탄사 활용",
                "'-요' 체의 친근한 말투",
                "이모티콘이나 감정을 나타내는 표현 사용",
            ],
        },
    }

    style_info = instructions.get(
        style, {"instruction": "자연스럽게 번역하세요.", "rules": []}
    )

    return (
        style_info["instruction"]
        + "\n"
        + "\n".join(f"- {rule}" for rule in style_info["rules"])
    )


def translate_text_using_chatgpt(text, src_lang, trg_lang, style):
    style_instruction = get_style_instruction(style)
    example = parallel_example.get(trg_lang, {}).get(style, [""])[0]

    messages = [
        {
            "role": "system",
            "content": f"당신은 전문 번역가입니다. {src_lang}을 {trg_lang}으로 번역합니다.\n\n{style_instruction}",
        },
        {"role": "user", "content": f"다음은 {style} 스타일의 예시입니다:\n{example}"},
        {"role": "user", "content": f"다음 텍스트를 위 스타일로 번역해주세요:\n{text}"},
    ]

    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo", messages=messages, temperature=0.7
    )

    return response.choices[0].message.content


# Streamlit UI
st.title("번역 서비스")

# 텍스트 입력
text = st.text_area("번역할 텍스트를 입력하세요", "")

# 파일 업로더 추가
uploaded_file = st.file_uploader(
    "또는 파일을 업로드하세요",
    type=["pdf", "docx", "pptx", "xlsx"],
    accept_multiple_files=False,
)

# 언어 선택
src_lang = st.selectbox("원본 언어", ["영어", "한국어", "중국어"])
trg_lang = st.selectbox("목표 언어", ["영어", "한국어", "중국어"], index=1)

# 번역 스타일 선택
style = st.selectbox(
    "번역 스타일",
    ["Academic", "Fluent", "Formal", "Creative", "Informal"],
    format_func=lambda x: (
        f"{x} - 학술적"
        if x == "Academic"
        else (
            f"{x} - 자연스러운"
            if x == "Fluent"
            else (
                f"{x} - 격식있는"
                if x == "Formal"
                else f"{x} - 창의적인" if x == "Creative" else f"{x} - 비격식"
            )
        )
    ),
)

if st.button("번역"):
    if uploaded_file is not None:
        # 파일 번역
        translated_file = translate_file(uploaded_file, src_lang, trg_lang, style)

        # 파일 이름 및 MIME 타입 설정
        if uploaded_file.type == "application/pdf":
            download_filename = uploaded_file.name.replace(".pdf", ".docx")
            mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        else:
            download_filename = f"translated_{uploaded_file.name}"
            mime_type = uploaded_file.type

        # 다운로드 버튼 생성
        st.download_button(
            label="번역된 파일 다운로드",
            data=translated_file,
            file_name=download_filename,
            mime=mime_type,
        )
    elif text:
        # 텍스트 번역
        translated_text = translate_text_using_chatgpt(text, src_lang, trg_lang, style)

        # 번역 결과를 컨테이너에 표시
        with st.container():
            # 두 개의 컬럼 생성
            col1, col2 = st.columns([6, 1])

            # 첫 번째 컬럼에 번역 결과 표시
            with col1:
                st.text_area("번역 결과", translated_text, height=200, key="result")

            # 두 번째 컬럼에 복사 버튼 표시
            with col2:
                if st.button("복사", key="copy"):
                    try:
                        import pyperclip

                        pyperclip.copy(translated_text)
                        st.success("복사 완료!")
                    except:
                        st.error("복사 실패")
    else:
        st.warning("번역할 텍스트를 입력하거나 파일을 업로드해주세요.")
